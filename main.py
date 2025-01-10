import os
import time
import logging
from sys import exit
from pathlib import Path
from contextlib import contextmanager
from selenium import webdriver
from selenium.webdriver.chrome.options import Options as ChromeOptions
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.action_chains import ActionChains
from selenium.common.exceptions import TimeoutException
from pptx import Presentation
from pptx.util import Cm
import win32com.client

# Get the user's Desktop path
desktop_path = os.path.join(os.path.join(os.environ['USERPROFILE']), 'Desktop')

# Configuration
CONFIG = {
    'url': "https://freemeteo.gr/kairos/plati/7-imeres/pinakas/?gid=734573&language=greek&country=greece",
    'screenshot_path': "weather_screenshot.png",
    'template_path': "template.pptx",
    'output_pptx': "updated_presentation.pptx",
    'output_video': os.path.join(desktop_path, "Καιρός.mp4"),  # Save video to Desktop
    'slide_dimensions': {'left': 4.18, 'top': 1.29, 'width': 17.03, 'height': 11.69}
}

# Setup logging with file handler and format
logging.basicConfig(
    level=logging.ERROR,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[ 
        logging.StreamHandler(),
        logging.FileHandler("app.log", mode='w', encoding='utf-8')
    ]
)

@contextmanager
def chrome_driver():
    """Context manager for Chrome driver"""
    options = ChromeOptions()
    options.add_argument("--start-maximized")
    options.add_argument("--headless")  # Add this line to run Chrome in headless mode
    options.add_argument("--disable-gpu")  # Recommended for headless mode
    options.add_argument("--no-sandbox")  # Optional: can improve stability in some environments
    options.add_argument("--disable-dev-shm-usage")  # Optional: avoid shared memory issues
    logging.debug("Initializing Chrome driver")
    driver = webdriver.Chrome(options=options)
    try:
        yield driver
    finally:
        logging.debug("Quitting Chrome driver")
        driver.quit()

def capture_weather_screenshot():
    """Capture weather data screenshot using Selenium"""
    with chrome_driver() as driver:
        try:
            logging.info("Accessing weather website")
            driver.get(CONFIG['url'])

            # Check and accept cookies if the button is present
            try:
                WebDriverWait(driver, 5).until(
                    EC.element_to_be_clickable((By.XPATH, "//*[contains(text(), 'Συναίνεση')]"))
                ).click()
                logging.debug("Cookies accepted")
            except TimeoutException:
                logging.debug("No cookie consent button found, proceeding without clicking")

            # Wait for and find weather element
            weather = WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.CLASS_NAME, "today.table"))
            )
            logging.debug("Weather data loaded successfully")

            # Move to tooltip element
            tooltip_element = driver.find_element(By.CLASS_NAME, "prev.sevendays")
            ActionChains(driver).move_to_element(tooltip_element).perform()
            logging.debug("Tooltip element interacted with")

            # Take screenshot
            weather.screenshot(CONFIG['screenshot_path'])
            logging.info("Weather screenshot captured successfully")
        except Exception as e:
            logging.error(f"Error during screenshot capture: {str(e)}")
            raise

def update_powerpoint():
    """Update PowerPoint presentation with screenshot"""
    try:
        logging.info("Opening PowerPoint template")
        prs = Presentation(CONFIG['template_path'])
        slide = prs.slides[0]
        dims = CONFIG['slide_dimensions']
        
        slide.shapes.add_picture(
            CONFIG['screenshot_path'],
            Cm(dims['left']),
            Cm(dims['top']),
            Cm(dims['width']),
            Cm(dims['height'])
        )
        prs.save(CONFIG['output_pptx'])
        logging.info("PowerPoint presentation updated successfully")
    except Exception as e:
        logging.error(f"Error updating PowerPoint: {str(e)}")
        raise

def create_video():
    """Convert PowerPoint to video"""
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint.DisplayAlerts = 0  # Suppress alerts

    try:
        logging.info("Opening PowerPoint file for video conversion")
        presentation = powerpoint.Presentations.Open(
            FileName=str(Path.cwd() / CONFIG['output_pptx']),
            WithWindow=True  # Must be True to avoid the "Visible" error
        )
        
        # Minimize the PowerPoint window
        powerpoint.WindowState = 2  # ppWindowMinimized
        
        presentation.CreateVideo(
            str(Path.cwd() / CONFIG['output_video']),
            VertResolution=1080,
            Quality=100
        )
        
        logging.debug("Waiting for video creation to complete")
        while presentation.CreateVideoStatus == 1:  # Check if video rendering is complete
            time.sleep(1)
        
        presentation.Close()
        logging.info("Video created successfully")
    except Exception as e:
        logging.error(f"Error creating video: {str(e)}")
        raise
    finally:
        powerpoint.Quit()  # Ensure PowerPoint quits


def main():
    """Main execution function"""
    try:
        logging.info("Starting process")
        capture_weather_screenshot()
        update_powerpoint()
        create_video()
        logging.info("Process completed successfully")
    except Exception as e:
        logging.error(f"Process failed: {str(e)}")
        return 1
    return 0

if __name__ == "__main__":
    main()
