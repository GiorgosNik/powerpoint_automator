import pytest
import os
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from main import chrome_driver, capture_weather_screenshot, update_powerpoint, CONFIG
from pptx import Presentation

@pytest.fixture
def setup_teardown():
    # Setup
    original_screenshot_path = CONFIG['screenshot_path']
    original_template_path = CONFIG['template_path']
    original_output_pptx = CONFIG['output_pptx']
    
    # Create test paths
    CONFIG['screenshot_path'] = 'test_screenshot.png'
    CONFIG['template_path'] = 'test_template.pptx'
    CONFIG['output_pptx'] = 'test_output.pptx'
    
    yield
    
    # Cleanup
    for file in [CONFIG['screenshot_path'], CONFIG['output_pptx']]:
        if os.path.exists(file):
            os.remove(file)
            
    # Restore original config
    CONFIG['screenshot_path'] = original_screenshot_path
    CONFIG['template_path'] = original_template_path  
    CONFIG['output_pptx'] = original_output_pptx

def test_chrome_driver():
    with chrome_driver() as driver:
        assert driver is not None
        assert driver.capabilities['browserName'].lower() == 'chrome'

@pytest.mark.flaky(reruns=2)
def test_capture_weather_screenshot(setup_teardown):
    try:
        capture_weather_screenshot()
        assert os.path.exists(CONFIG['screenshot_path'])
        assert os.path.getsize(CONFIG['screenshot_path']) > 0
    except Exception as e:
        pytest.fail(f"Screenshot capture failed: {str(e)}")

def test_update_powerpoint(setup_teardown, tmp_path):
    # Create a dummy screenshot file
    from PIL import Image
    dummy_image = Image.new('RGB', (100, 100), color='red')
    dummy_screenshot = tmp_path / "test_screenshot.png"
    dummy_image.save(dummy_screenshot)
    CONFIG['screenshot_path'] = str(dummy_screenshot)
    
    # Create a dummy PowerPoint template file
    dummy_template = tmp_path / "test_template.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(dummy_template)
    CONFIG['template_path'] = str(dummy_template)
    
    try:
        update_powerpoint()
        assert os.path.exists(CONFIG['output_pptx'])
        assert os.path.getsize(CONFIG['output_pptx']) > 0
    except Exception as e:
        pytest.fail(f"PowerPoint update failed: {str(e)}")

def test_weather_website_accessibility():
    with chrome_driver() as driver:
        try:
            driver.get(CONFIG['url'])
            WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.CLASS_NAME, "today.table"))
            )
            assert "freemeteo.gr" in driver.current_url
        except Exception as e:
            pytest.fail(f"Weather website not accessible: {str(e)}")

def test_cookie_consent_button():
    with chrome_driver() as driver:
        try:
            driver.get(CONFIG['url'])
            consent_button = WebDriverWait(driver, 10).until(
                EC.element_to_be_clickable((By.XPATH, "//*[contains(text(), 'Συναίνεση')]"))
            )
            assert consent_button.is_displayed()
            assert consent_button.is_enabled()
        except Exception as e:
            pytest.fail(f"Cookie consent button not found: {str(e)}")

def test_update_powerpoint_invalid_template():
    invalid_template = "nonexistent_template.pptx"
    original_template = CONFIG['template_path']
    CONFIG['template_path'] = invalid_template
    
    with pytest.raises(Exception):
        update_powerpoint()
        
    CONFIG['template_path'] = original_template

def test_update_powerpoint_invalid_screenshot():
    invalid_screenshot = "nonexistent_screenshot.png"
    original_screenshot = CONFIG['screenshot_path']
    CONFIG['screenshot_path'] = invalid_screenshot
    
    with pytest.raises(Exception):
        update_powerpoint()
        
    CONFIG['screenshot_path'] = original_screenshot